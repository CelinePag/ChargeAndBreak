# -*- coding: utf-8 -*-
"""
Editable PowerPoint version of the conceptual solution figure.

Same data reconstruction as concept_solution_figure.py, but every element is
a native PowerPoint shape (rectangles, triangles, connectors, freeform
curves, text boxes) so colours / positions / text can be adjusted directly
in PowerPoint.  Shapes are named and grouped per panel (Route / Callouts /
Truck / SOC / HoS / Driver) -- use the Selection Pane (Alt+F10) to navigate.

Run:  python -m src.plot.concept_solution_pptx
"""
import json
import math
import os

from src import paths as _paths

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE, MSO_PATTERN
from pptx.dml.color import RGBColor

BASE = str(_paths.ROOT)

# --------------------------------------------------------------------------
# Per-run layout config.  Callout / micro-label TEXTS are generated from the
# solution automatically; only their x/y positions are tuned per run.
#   callout_pos : stop -> (text_center_x [h], row)   row 0 = top, 1 = lower
#   micro       : list of (stop, block_key, label_x [h], label_y [in])
#   source      : "sim" (executed plan, default) | "oracle" (hindsight optimum
#                 stored in the same solution file, key sol_file)
#   soc_note    : optional ((lines), text_x [h], text_y [kWh], tip_x, tip_y)
# --------------------------------------------------------------------------
RUNS = {
    "RshortCfewTmedium_19_RO_20260716_092310_060": dict(
        inst="RshortCfewTmedium_19",
        callout_pos={10: (8.15, 0), 11: (12.85, 0), 20: (20.9, 0),
                     27: (27.1, 0), 35: (31.0, 0)},
        micro=[(10, "mstop", 9.70, 2.79), (10, "queue", 10.60, 2.61),
               (10, "charge", 11.30, 2.82), (20, "mseq", 16.60, 2.66)],
    ),
    "RshortCfewTlarge_24_RO_20260716_092310_016": dict(
        inst="RshortCfewTlarge_24",
        callout_pos={6: (8.10, 0), 9: (12.10, 1), 17: (16.10, 0),
                     22: (22.30, 0), 27: (27.90, 0), 36: (32.30, 1)},
        micro=[(17, "mstop", 13.60, 2.82), (17, "queue", 14.15, 2.60),
               (17, "charge", 14.75, 2.84)],
    ),
    "RshortCfewTlarge_24_ORACLE": dict(
        inst="RshortCfewTlarge_24",
        sol_file="RshortCfewTlarge_24_RO_20260716_092310_016",
        source="oracle",
        callout_pos={12: (13.30, 0), 22: (21.00, 0), 24: (27.00, 0),
                     36: (31.30, 1)},
        micro=[(12, "mstop", 11.30, 2.82), (12, "queue", 11.85, 2.60),
               (12, "charge", 12.35, 2.84)],
        soc_note=([[("arrives exactly at E", 0), ("min", 30)],
                   "(no safety buffer)"],
                  24.9, 250.0, 26.40, 102.0),
    ),
}

import sys
SOL  = sys.argv[1] if len(sys.argv) > 1 else "RshortCfewTlarge_24_ORACLE"
CFG  = RUNS[SOL]
INST = CFG["inst"]
SOLFILE = CFG.get("sol_file", SOL)
SOURCE  = CFG.get("source", "sim")
OUT  = _paths.figures(f"solution_concept_{SOL}.pptx")

# ================================================================ data ======
sol  = json.load(open(_paths.solutions(SOLFILE + ".json")))
inst = json.load(open(_paths.instances(INST + ".json")))
fd     = inst["instance"]
D_real = inst["D_real"]
E_real = inst["E_real"]
traj   = sol["sim_trajectory"]
acts   = sol["actions"]

N  = fd["N"]
K  = set(fd["K"]); C = set(fd["C"]); L = set(fd["L"])
km = fd["km"]
cumkm = [0.0]
for i in range(N):
    cumkm.append(cumkm[-1] + km[str(i)])
TOTKM = cumkm[-1]

Ebar = {int(k): v for k, v in fd["Ebar"].items()}
Tbar = {int(k): v for k, v in fd["Tbar"].items()}
rs   = sorted(Ebar); Es = [Ebar[r] for r in rs]; Ts = [Tbar[r] for r in rs]

def e2t(e):
    e = max(Es[0], min(Es[-1], e))
    for k in range(len(Es) - 1):
        if Es[k] <= e <= Es[k + 1]:
            span = Es[k + 1] - Es[k]
            return Ts[k] + (e - Es[k]) / span * (Ts[k + 1] - Ts[k]) if span else Ts[k]
    return Ts[-1]

events, soc_pts, drive_blocks, charge_bands = [], [], [], []
cd_resets, rest_resets = [], []

def _soc_charge_path(t_start, e_arr, e_dep):
    """append SOC points along the PWL charging curve, incl. knot crossings"""
    t0c = e2t(e_arr)
    soc_pts.append((t_start, e_arr))
    for Ek, Tk in zip(Es, Ts):
        if e_arr < Ek < e_dep:
            soc_pts.append((t_start + (Tk - t0c), Ek))
    soc_pts.append((t_start + (e2t(e_dep) - t0c), e_dep))

if SOURCE == "oracle":
    # ---- hindsight-optimal plan: embedded (old runs) or shared cache --------
    _oracle = sol.get("oracle")
    if not _oracle:
        from src.methods.oracle import load_oracle_cache
        _oracle = load_oracle_cache(sol.get("instance", "")) or {}
    osol = _oracle["sol"]
    soc_pts.append((osol[0]["ta"], osol[0]["ea"]))
    for k in range(len(osol) - 1):
        s, s1 = osol[k], osol[k + 1]
        i  = s["i"]
        ta, td = s["ta"], s["td"]
        dwell = td - ta
        brk = ("b45" if s["b45"] else "b15" if s["b15"] else
               "b30" if s["b30"] else None)
        rst = "r1" if s["rho1"] else "r2" if s["rho2"] else None
        blocks = []
        tt = ta
        if dwell > 1e-4:
            if i in K:
                v = 1 if (s["y"] or s["tauc"] > 1e-6 or s["taub"] > 1e-6
                          or s["taur"] > 1e-6) else 0
                if v:
                    ms = fd["M_stop"][str(i)]
                    blocks.append(("mstop", tt, tt + ms)); tt += ms
                if s["tauq"] > 1e-6:
                    blocks.append(("queue", tt, tt + s["tauq"])); tt += s["tauq"]
                if s["tauc"] > 1e-6:
                    blocks.append(("charge", tt, tt + s["tauc"]))
                    charge_bands.append((tt, tt + s["tauc"]))
                    _soc_charge_path(tt, s["ea"], s["ed"])
                    tt += s["tauc"]
                if s["sigma"]:
                    mq = fd["M_seq"][str(i)]
                    blocks.append(("mseq", tt, tt + mq)); tt += mq
                if s["taub"] > 1e-6:
                    blocks.append(("break", tt, tt + s["taub"])); tt += s["taub"]
                if s["taur"] > 1e-6:
                    blocks.append(("rest", tt, tt + s["taur"])); tt += s["taur"]
            elif i in C:
                sv = fd["S"][str(i)]
                blocks.append(("service", tt, tt + sv)); tt += sv
                if s["taur"] > 1e-6:
                    blocks.append(("rest", tt, tt + s["taur"])); tt += s["taur"]
                elif s["taub"] > 1e-6:
                    blocks.append(("break", tt, tt + s["taub"])); tt += s["taub"]
            elif i in L:
                if s["taub"] > 1e-6 or s["taur"] > 1e-6:
                    ml = fd["M_lay"][str(i)]
                    blocks.append(("mlay", tt, tt + ml)); tt += ml
                if s["taub"] > 1e-6:
                    blocks.append(("break", tt, tt + s["taub"])); tt += s["taub"]
                if s["taur"] > 1e-6:
                    blocks.append(("rest", tt, tt + s["taur"])); tt += s["taur"]
            if td - tt > 1e-3:                     # residual = declared wait
                blocks.append(("idle", tt, tt + (td - tt)))
            events.append(dict(stop=i, t_arr=ta, td=td, blocks=blocks,
                               y=s["y"], brk=brk, rst=rst, sigma=s["sigma"]))
        if brk in ("b45", "b30") or rst:
            cd_resets.append(td)
        if rst:
            rest_resets.append(td)
        drive_blocks.append((td, s1["ta"]))
        soc_pts.append((td, s["ed"]))
        soc_pts.append((s1["ta"], s1["ea"]))
    T0, TEND = osol[0]["ta"], osol[-1]["ta"]
    TW_OK = all(not s["delta"] for s in osol if s["is_C"])
    DUR_H = TEND - T0
else:
    # ---- executed plan reconstructed from the simulation trajectory --------
    soc_pts.append((traj[0]["t_arr"], traj[0]["e_arr"]))
    for i in range(N):
        s, s1 = traj[i], traj[i + 1]
        a  = acts[i]
        td = s1["t_arr"] - D_real[i]
        dwell = td - s["t_arr"]
        y, brk, rst = a.get("y"), a.get("break_type"), a.get("rest_type")
        e_arr = s["e_arr"]
        e_dep = s1["e_arr"] + E_real[i]
        chg   = e_dep - e_arr
        blocks = []
        tt = s["t_arr"]
        if dwell > 1e-3:
            sigma = 0
            if i in K:
                ms = fd["M_stop"][str(i)]
                qv = fd["Q"][str(i)] if y else 0.0
                tauc = e2t(e_dep) - e2t(e_arr) if chg > 1e-6 else 0.0
                sigma = 1 if (y and rst) else 0
                mseq = fd["M_seq"][str(i)] * sigma
                taur = fd["Tr2"] if rst == "r2" else (fd["Tr1"] if rst == "r1" else 0.0)
                extra = dwell - (ms + qv + tauc + mseq + taur)
                blocks.append(("mstop", tt, tt + ms)); tt += ms
                if qv > 1e-6:
                    blocks.append(("queue", tt, tt + qv)); tt += qv
                if tauc > 1e-6:
                    blocks.append(("charge", tt, tt + tauc))
                    charge_bands.append((tt, tt + tauc))
                    _soc_charge_path(tt, e_arr, e_dep)
                    tt += tauc
                if mseq > 1e-6:
                    blocks.append(("mseq", tt, tt + mseq)); tt += mseq
                if brk and extra > 1e-6:
                    blocks.append(("break", tt, tt + extra)); tt += extra
                if taur > 1e-6:
                    blocks.append(("rest", tt, tt + taur)); tt += taur
            elif i in C:
                sv = fd["S"][str(i)]
                blocks.append(("service", tt, tt + sv)); tt += sv
                rem = dwell - sv
                if rst and rem > 1e-6:
                    blocks.append(("rest", tt, tt + rem)); tt += rem
                elif brk and rem > 1e-6:
                    blocks.append(("break", tt, tt + rem)); tt += rem
            elif i in L:
                ml = fd["M_lay"][str(i)] if (brk or rst) else 0.0
                if ml > 1e-6:
                    blocks.append(("mlay", tt, tt + ml)); tt += ml
                rem = td - tt
                kind = "rest" if rst else ("break" if brk else "idle")
                if rem > 1e-6:
                    blocks.append((kind, tt, tt + rem)); tt += rem
            events.append(dict(stop=i, t_arr=s["t_arr"], td=td, blocks=blocks,
                               y=y, brk=brk, rst=rst, sigma=sigma))
        if brk in ("b45", "b30") or rst:
            cd_resets.append(td)
        if rst:
            rest_resets.append(td)
        drive_blocks.append((td, s1["t_arr"]))
        soc_pts.append((td, e_dep))
        soc_pts.append((s1["t_arr"], s1["e_arr"]))
    T0, TEND = traj[0]["t_arr"], traj[-1]["t_arr"]
    TW_OK = (sol["metrics"]["tw_n_misses"] == 0)
    DUR_H = sol["duration_h"]

_md = []
for (t0, t1) in drive_blocks:
    if _md and abs(t0 - _md[-1][1]) < 2.5e-3:
        _md[-1][1] = t1
    else:
        _md.append([t0, t1])
drive_blocks = [(a, b) for a, b in _md]

def counter_path(reset_times):
    pts, val = [(T0, 0.0)], 0.0
    segs = [("drive", a, b) for a, b in drive_blocks] + \
           [("reset", r, r) for r in reset_times]
    segs.sort(key=lambda x: (x[1], 0 if x[0] == "reset" else 1))
    cur = T0
    for kind, t0, t1 in segs:
        if t0 > cur:
            pts.append((t0, val))
        if kind == "reset":
            pts.append((t0, val)); val = 0.0; pts.append((t0, 0.0)); cur = t0
        else:
            val += (t1 - t0); pts.append((t1, val)); cur = t1
    pts.append((TEND, val))
    return pts

cd_path = counter_path(cd_resets)
sd_path = counter_path(rest_resets)

# ============================================================ styling =======
def RGB(h): return RGBColor.from_string(h.lstrip("#"))

INK, MUT   = "222222", "666666"
c_drive    = "D9D9D9"; c_drive_e = "9A9A9A"
c_mstop    = "141414"
c_queue    = "E02B2B"
c_charge   = "1E7A3C"; c_charge_dk = "0D3C1D"
c_cs_faded = "8FBFA0"
c_work     = "F2C7EE"; c_work_e  = "C583BE"
c_serv     = "E8822D"; c_serv_e  = "A85A17"
c_brk_f    = "EEF6FD"; c_brk_e   = "2B7BBA"
c_rst_f    = "E4EDF7"; c_rst_e   = "1F4E79"
c_ring     = "CC2222"
c_cd       = "2B6CB0"
c_sd       = "8A63B8"
FONT       = "Calibri"

# ============================================================ layout ========
SLIDE_W, SLIDE_H = 40 / 3.0, 7.5          # inches (exact 16:9)
PL, PR = 1.02, 12.90                      # shared panel left/right
PW = PR - PL
TMIN, TMAX = T0 - 0.35, TEND + 0.35
KMIN, KMAX = -18.0, TOTKM + 18.0

def xt(t):  return PL + (t - TMIN) / (TMAX - TMIN) * PW      # time -> x [in]
def xk(k):  return PL + (k - KMIN) / (KMAX - KMIN) * PW      # km   -> x [in]

ROUTE_Y  = 1.46                           # route baseline
CALL_Y   = 2.12                           # callout text top
TRUCK_Y, TRUCK_H = 3.02, 0.48
SOC_Y,   SOC_H   = 3.70, 0.90
HOS_Y,   HOS_H   = 4.78, 0.84
DRV_Y,   DRV_H   = 5.95, 0.48
EMAX_SOC, HMAX_HOS = 560.0, 10.3

def ye(e): return SOC_Y + SOC_H - e / EMAX_SOC * SOC_H       # kWh -> y
def yh(v): return HOS_Y + HOS_H - v / HMAX_HOS * HOS_H       # h   -> y

# ============================================================ helpers =======
prs = Presentation()
prs.slide_width  = Emu(12192000)
prs.slide_height = Inches(SLIDE_H)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

def box(sh, x, y, w, h, fill=None, line=None, lw=0.75, name=None,
        shape=MSO_SHAPE.RECTANGLE, pattern=None, back=None, dash=None):
    s = sh.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if pattern is not None:
        s.fill.patterned()
        s.fill.pattern = pattern
        s.fill.fore_color.rgb = RGB(fill)
        s.fill.back_color.rgb = RGB(back or "FFFFFF")
    elif fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = RGB(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = RGB(line)
        s.line.width = Pt(lw)
        if dash is not None:
            s.line.dash_style = dash
    if name:
        s.name = name
    return s

def conn(sh, x1, y1, x2, y2, color, lw=1.0, dash=None, name=None):
    c = sh.add_connector(MSO_CONNECTOR.STRAIGHT,
                         Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.shadow.inherit = False
    c.line.color.rgb = RGB(color)
    c.line.width = Pt(lw)
    if dash is not None:
        c.line.dash_style = dash
    if name:
        c.name = name
    return c

def _fmt_run(r, size, color, bold, italic):
    r.font.size = Pt(size)
    r.font.name = FONT
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGB(color)

def txt(sh, x, y, w, h, lines, size=10, color=INK, bold=False, italic=False,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, rotation=0.0, name=None):
    """lines: str | list of lines; a line is a str or list of (text, baseline%) runs"""
    tb = sh.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for li, ln in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(ln, str):
            ln = [(ln, 0)]
        for text, basepct in ln:
            r = p.add_run()
            r.text = text
            _fmt_run(r, size, color, bold, italic)
            if basepct:
                r.font._rPr.set("baseline", str(int(basepct * 1000)))
    if rotation:
        tb.rotation = rotation
    if name:
        tb.name = name
    return tb

def polyline(sh, pts, color, lw=1.75, dash=None, name=None):
    """pts: list of (x,y) in inches -> open freeform polyline"""
    e = [(int(Inches(x)), int(Inches(y))) for x, y in pts]
    fb = sh.build_freeform(e[0][0], e[0][1], scale=1.0)
    fb.add_line_segments(e[1:], close=False)
    s = fb.convert_to_shape()
    s.shadow.inherit = False
    s.fill.background()
    s.line.color.rgb = RGB(color)
    s.line.width = Pt(lw)
    if dash is not None:
        s.line.dash_style = dash
    if name:
        s.name = name
    return s

# math-ish label parts
M_STOP = [("M", 0), ("i", -25), ("stop", 30)]
M_SEQ  = [("M", 0), ("i", -25), ("seq", 30)]
Q_I    = [("Q", 0), ("i", -25)]
TAU_C  = [("τ", 0), ("i", -25), ("c", 30)]
TAU_R  = [("τ", 0), ("i", -25), ("r", 30)]

# ============================================================ title =========
arr_clock = f"{int(TEND % 24):02d}:{int(round((TEND % 1) * 60)):02d}"
txt(slide.shapes, 0.85, 0.18, 6.0, 0.45, "What a solution looks like",
    size=24, bold=True, align=PP_ALIGN.LEFT, name="title")
cust_word = "customer" if len(C) == 1 else "customers"
tw_txt = (("customer window met" if len(C) == 1 else
           f"all {len(C)} customer windows met") if TW_OK else
          "some customer windows missed")
src_txt = ("Oracle plan (hindsight optimum, solved with the realised travel "
           "times)" if SOURCE == "oracle" else "RO plan, executed as-is")
txt(slide.shapes, 0.87, 0.62, 8.5, 0.24,
    f"Instance {INST}  ({TOTKM:.0f} km, {len(K)} charging stations, "
    f"{len(L)} rest areas, {len(C)} {cust_word})",
    size=11, color=MUT, align=PP_ALIGN.LEFT, name="subtitle.1")
txt(slide.shapes, 0.87, 0.86, 11.0, 0.24,
    f"{src_txt}  •  depart 08:00, arrive {arr_clock} next day "
    f"({DUR_H:.1f} h)  •  feasible, {tw_txt}",
    size=11, color=MUT, align=PP_ALIGN.LEFT, name="subtitle.2")

# ============================================================ route =========
g = slide.shapes.add_group_shape(); g.name = "panel.route"
gs = g.shapes
txt(gs, 0.10, ROUTE_Y - 0.12, 0.8, 0.25, "Route", size=13, bold=True,
    align=PP_ALIGN.LEFT, name="route.label")
conn(gs, xk(0), ROUTE_Y, xk(TOTKM), ROUTE_Y, "444444", 1.5, name="route.line")

chosen = {e["stop"] for e in events}
for i in range(N + 1):
    x = xk(cumkm[i])
    if i in L:
        box(gs, x - 0.008, ROUTE_Y - 0.06, 0.016, 0.12, fill="999999",
            name=f"route.layby.{i}")
    if i in K:
        sel = i in chosen
        w = 0.21 if sel else 0.155
        h = 0.18 if sel else 0.135
        box(gs, x - w / 2, ROUTE_Y - 0.035 - h, w, h,
            fill=(c_charge if sel else c_cs_faded),
            line=(c_charge_dk if sel else "7AA88C"), lw=0.75,
            shape=MSO_SHAPE.ISOSCELES_TRIANGLE,
            name=f"route.cs.{i}" + (".chosen" if sel else ""))
    if i in C:
        d = 0.18
        box(gs, x - d / 2, ROUTE_Y - 0.035 - d, d, d, fill=c_serv,
            line="7A3D0D", lw=0.75, shape=MSO_SHAPE.DIAMOND,
            name=f"route.customer.{i}")
for i in sorted(chosen):
    x = xk(cumkm[i])
    r = 0.34
    box(gs, x - r / 2, ROUTE_Y - 0.125 - r / 2, r, r, fill=None,
        line=c_ring, lw=1.75, shape=MSO_SHAPE.OVAL, name=f"route.ring.{i}")
prev_x, row2 = -1e9, False
for i in sorted(chosen):
    x = xk(cumkm[i])
    row2 = (cumkm[i] - prev_x < 45) and not row2
    txt(gs, x - 0.45, ROUTE_Y + (0.34 if row2 else 0.14), 0.9, 0.2,
        f"{cumkm[i]:.0f} km", size=9, color=c_ring, bold=True,
        name=f"route.km.{i}")
    prev_x = cumkm[i]
for x_km, lab in [(0.0, "O"), (TOTKM, "D")]:
    x = xk(x_km)
    box(gs, x - 0.09, ROUTE_Y - 0.09, 0.18, 0.18, fill="111111",
        name=f"route.node.{lab}")
    txt(gs, x - 0.3, ROUTE_Y - 0.42, 0.6, 0.24, lab, size=12, bold=True,
        name=f"route.nodelabel.{lab}")
    txt(gs, x - 0.45, ROUTE_Y + 0.14, 0.9, 0.2, f"{x_km:.0f} km", size=9,
        name=f"route.kmlabel.{lab}")
for ci, c in enumerate(sorted(C), start=1):
    cx = xk(cumkm[c])
    txt(gs, cx + 0.10, ROUTE_Y - 0.44, 0.5, 0.22, f"C{ci}", size=10, bold=True,
        color="7A3D0D", align=PP_ALIGN.LEFT, name=f"route.C{ci}")

# legend
lg = slide.shapes.add_group_shape(); lg.name = "legend"
lgs = lg.shapes
LX, LY = 7.02, 0.36
box(lgs, LX, LY, 5.92, 0.36, fill="FFFFFF", line="33415C", lw=1.0,
    name="legend.frame")
box(lgs, LX + 0.18, LY + 0.10, 0.016, 0.16, fill="999999", name="legend.layby")
txt(lgs, LX + 0.26, LY + 0.085, 0.85, 0.2, "Rest area", size=9.5,
    align=PP_ALIGN.LEFT, name="legend.layby.txt")
box(lgs, LX + 1.22, LY + 0.09, 0.17, 0.15, fill=c_charge, line=c_charge_dk,
    lw=0.75, shape=MSO_SHAPE.ISOSCELES_TRIANGLE, name="legend.cs")
txt(lgs, LX + 1.46, LY + 0.085, 1.45, 0.2, "Charging station", size=9.5,
    align=PP_ALIGN.LEFT, name="legend.cs.txt")
box(lgs, LX + 3.02, LY + 0.09, 0.16, 0.16, fill=c_serv, line="7A3D0D",
    lw=0.75, shape=MSO_SHAPE.DIAMOND, name="legend.cust")
txt(lgs, LX + 3.25, LY + 0.085, 0.9, 0.2, "Customer", size=9.5,
    align=PP_ALIGN.LEFT, name="legend.cust.txt")
box(lgs, LX + 4.18, LY + 0.075, 0.2, 0.2, fill=None, line=c_ring, lw=1.5,
    shape=MSO_SHAPE.OVAL, name="legend.ring")
txt(lgs, LX + 4.45, LY + 0.085, 1.45, 0.2, "Stop used by the solution",
    size=9.5, align=PP_ALIGN.LEFT, name="legend.ring.txt")

# ========================================================= callouts =========
g = slide.shapes.add_group_shape(); g.name = "panel.callouts"
gs = g.shapes
CALLW = 2.6
G_EQ_TAUC = [("(g", 0), ("i", -25), (" = τ", 0), ("i", -25), ("c", 30), (")", 0)]

def callout_lines(e, first_b45):
    """Generate callout text from the event structure."""
    i = e["stop"]
    has_chg = any(k == "charge" for k, _, _ in e["blocks"])
    rst_h = {"r1": "11 h", "r2": "9 h"}.get(e["rst"])
    if i in K:
        if e["rst"]:
            if has_chg:
                return ["Charge, then %s rest" % rst_h,
                        [("(sequential:  σ", 0), ("i", -25), (" = 1,  g", 0),
                         ("i", -25), (" = 0)", 0)]]
            return ["%s rest at charging station" % rst_h]
        if e["brk"] == "b45":
            return ["Charge & 45 min break"] + ([G_EQ_TAUC] if first_b45 else [])
        if e["brk"] == "b15":
            return ["Charge & 15 min break",
                    [("(1st split,  g", 0), ("i", -25), (" = τ", 0),
                     ("i", -25), ("c", 30), (")", 0)]]
        if e["brk"] == "b30":
            return ["Charge & 30 min break", "(2nd split)"]
        return ["Charge"]
    if i in C:
        if e["rst"]:
            return ["Service, then %s rest" % rst_h, "(overnight at customer)"]
        if e["brk"] == "b30":
            return ["Service + 30 min break", "(2nd split)"]
        if e["brk"] == "b15":
            return ["Service + 15 min break", "(1st split)"]
        return None                      # plain service: no callout
    if e["rst"]:                         # layby
        return ["%s rest at rest area" % rst_h]
    return [{"b45": "45 min break", "b15": "15 min break (1st split)",
             "b30": "30 min break (2nd split)"}[e["brk"]], "at rest area"]

seen_b45 = False
for e in events:
    i = e["stop"]
    lines = callout_lines(e, first_b45=(e["brk"] == "b45" and not seen_b45))
    if e["brk"] == "b45":
        seen_b45 = True
    if lines is None:
        continue
    if i not in CFG["callout_pos"]:
        print(f"NOTE: stop {i} has a callout but no position configured")
        continue
    tx, row = CFG["callout_pos"][i]
    y0 = CALL_Y + row * 0.42
    xm = 0.5 * (e["t_arr"] + e["td"])
    txt(gs, xt(tx) - CALLW / 2, y0, CALLW, 0.40, lines, size=9.5,
        color="444444", name=f"callout.{i}.txt")
    y_txt_bot = y0 + 0.185 * len(lines) + 0.03
    conn(gs, xt(tx), y_txt_bot, xt(xm), TRUCK_Y - 0.03, "B8B8B8", 0.75,
         name=f"callout.{i}.leader")

# micro labels (positions from CFG)
MICRO_PARTS = dict(mstop=M_STOP, mseq=M_SEQ, queue=Q_I, charge=TAU_C,
                   mlay=[("M", 0), ("i", -25), ("lay", 30)])
for stop, key, tx, ty in CFG["micro"]:
    ev = next(e for e in events if e["stop"] == stop)
    bl = {k: (a, b) for k, a, b in ev["blocks"]}
    t0, t1 = bl[key]
    txt(gs, xt(tx) - 0.3, ty - 0.10, 0.6, 0.2, [MICRO_PARTS[key]], size=10,
        italic=True, name=f"micro.{key}.{stop}")
    conn(gs, xt(tx), ty + 0.09, xt(0.5 * (t0 + t1)), TRUCK_Y + 0.02,
         "555555", 0.75, name=f"micro.{key}.{stop}.leader")

# ============================================================ truck =========
g = slide.shapes.add_group_shape(); g.name = "panel.truck"
gs = g.shapes
txt(gs, 0.10, TRUCK_Y + TRUCK_H / 2 - 0.12, 0.8, 0.25, "Truck", size=13,
    bold=True, align=PP_ALIGN.LEFT, name="truck.label")
box(gs, PL, TRUCK_Y, PW, TRUCK_H, fill=None, line="AAAAAA", lw=1.0,
    name="truck.frame")
BY, BH = TRUCK_Y + 0.034, TRUCK_H - 0.068
kindfill = dict(mstop=(c_mstop, c_mstop), mseq=(c_mstop, c_mstop),
                queue=(c_queue, c_queue), charge=(c_charge, c_charge))
for (t0, t1) in drive_blocks:
    s = box(gs, xt(t0), BY, xt(t1) - xt(t0), BH, fill=c_drive,
            line=c_drive_e, lw=0.5, name=f"truck.drive@{t0:.1f}")
    if t1 - t0 > 1.55:
        s.text_frame.text = "Driving"
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _fmt_run(s.text_frame.paragraphs[0].runs[0], 9, "333333", False, False)
for e in events:
    white_run = None
    for kind, t0, t1 in e["blocks"] + [("_end", None, None)]:
        if kind in kindfill:
            if white_run:
                box(gs, xt(white_run[0]), BY, xt(white_run[1]) - xt(white_run[0]),
                    BH, fill="FFFFFF", line=c_drive_e, lw=0.5,
                    name=f"truck.idle@{white_run[0]:.1f}")
                white_run = None
            f, ec = kindfill[kind]
            box(gs, xt(t0), BY, xt(t1) - xt(t0), BH, fill=f, line=ec, lw=0.4,
                name=f"truck.{kind}@{t0:.1f}")
        elif kind == "_end":
            if white_run:
                box(gs, xt(white_run[0]), BY, xt(white_run[1]) - xt(white_run[0]),
                    BH, fill="FFFFFF", line=c_drive_e, lw=0.5,
                    name=f"truck.idle@{white_run[0]:.1f}")
        else:
            white_run = [white_run[0], t1] if white_run else [t0, t1]
rest_ev = next((e for e in events if e["rst"]), None)
if rest_ev is not None:
    t0, t1 = {k: (a, b) for k, a, b in rest_ev["blocks"]}["rest"]
    txt(gs, xt(0.5 * (t0 + t1)) - 1.0, TRUCK_Y + TRUCK_H / 2 - 0.10, 2.0, 0.2,
        [[("truck parked  (", 0)] + TAU_R + [(")", 0)]], size=9, italic=True,
        color="777777", name="truck.parked.txt")

# ============================================================== SOC =========
g = slide.shapes.add_group_shape(); g.name = "panel.soc"
gs = g.shapes
txt(gs, 0.10, SOC_Y + SOC_H / 2 - 0.22, 0.85, 0.45, ["SOC", "[kWh]"],
    size=13, bold=True, align=PP_ALIGN.LEFT, name="soc.label")
for (t0, t1) in charge_bands:
    box(gs, xt(t0), SOC_Y + 0.01, xt(t1) - xt(t0), SOC_H - 0.02,
        fill="DDEAE2", name=f"soc.chargeband@{t0:.1f}")
box(gs, PL, SOC_Y, PW, SOC_H, fill=None, line="AAAAAA", lw=1.0,
    name="soc.frame")
conn(gs, PL, ye(fd["Ecap"]), PR, ye(fd["Ecap"]), "888888", 1.0,
     dash=MSO_LINE.DASH, name="soc.Ecap.line")
conn(gs, PL, ye(fd["Emin"]), PR, ye(fd["Emin"]), "D62728", 1.0,
     dash=MSO_LINE.DASH, name="soc.Emin.line")
txt(gs, PR + 0.03, ye(fd["Ecap"]) - 0.09, 0.40, 0.18,
    [[("E", 0), ("cap", 30)]], size=9, color="666666",
    align=PP_ALIGN.LEFT, name="soc.Ecap.txt")
txt(gs, PR + 0.03, ye(fd["Emin"]) - 0.09, 0.40, 0.18,
    [[("E", 0), ("min", 30)]], size=9, color="D62728",
    align=PP_ALIGN.LEFT, name="soc.Emin.txt")
for v in (100, 300, 500):
    txt(gs, PL - 0.42, ye(v) - 0.09, 0.35, 0.18, str(v), size=8,
        color="444444", align=PP_ALIGN.RIGHT, name=f"soc.ytick.{v}")
polyline(gs, [(xt(t), ye(e)) for t, e in soc_pts], c_charge, 1.75,
         name="soc.curve")
conn(gs, xt(24), SOC_Y + 0.01, xt(24), SOC_Y + SOC_H - 0.01, "BBBBBB", 0.9,
     dash=MSO_LINE.ROUND_DOT, name="soc.midnight.line")
txt(gs, xt(24.15), ye(430) - 0.09, 0.75, 0.18, "midnight", size=8,
    color="999999", align=PP_ALIGN.LEFT, name="soc.midnight.txt")
if CFG.get("soc_note"):
    lines, nx, ny, tipx, tipy = CFG["soc_note"]
    txt(gs, xt(nx) - 0.9, ye(ny) - 0.10, 1.8, 0.34, lines, size=8,
        italic=True, color="D62728", name="soc.note.txt")
    conn(gs, xt(nx) + 0.45, ye(ny) + 0.26, xt(tipx) - 0.02, ye(tipy),
         "D62728", 0.75, name="soc.note.leader")

# ============================================================== HoS =========
g = slide.shapes.add_group_shape(); g.name = "panel.hos"
gs = g.shapes
txt(gs, 0.10, HOS_Y + HOS_H / 2 - 0.22, 0.85, 0.45, ["HoS", "[h]"],
    size=13, bold=True, align=PP_ALIGN.LEFT, name="hos.label")
box(gs, PL, HOS_Y, PW, HOS_H, fill=None, line="AAAAAA", lw=1.0,
    name="hos.frame")
conn(gs, PL, yh(fd["Tdrv_cons"]), PR, yh(fd["Tdrv_cons"]), c_cd, 1.0,
     dash=MSO_LINE.ROUND_DOT, name="hos.cap45.line")
conn(gs, PL, yh(fd["Tdrv_sh1"]), PR, yh(fd["Tdrv_sh1"]), c_sd, 1.0,
     dash=MSO_LINE.ROUND_DOT, name="hos.cap9.line")
txt(gs, PR + 0.03, yh(fd["Tdrv_cons"]) - 0.09, 0.42, 0.18, "4.5 h", size=8,
    color=c_cd, align=PP_ALIGN.LEFT, name="hos.cap45.txt")
txt(gs, PR + 0.03, yh(fd["Tdrv_sh1"]) - 0.09, 0.42, 0.18, "9 h", size=8,
    color=c_sd, align=PP_ALIGN.LEFT, name="hos.cap9.txt")
for v in (0, 4.5, 9):
    txt(gs, PL - 0.42, yh(v) - 0.09, 0.35, 0.18,
        f"{v:g}", size=8, color="444444", align=PP_ALIGN.RIGHT,
        name=f"hos.ytick.{v}")
conn(gs, xt(24), HOS_Y + 0.01, xt(24), HOS_Y + HOS_H - 0.01, "BBBBBB", 0.9,
     dash=MSO_LINE.ROUND_DOT, name="hos.midnight.line")
polyline(gs, [(xt(t), yh(v)) for t, v in cd_path], c_cd, 1.75,
         name="hos.cd.curve")
polyline(gs, [(xt(t), yh(v)) for t, v in sd_path], c_sd, 1.5,
         dash=MSO_LINE.DASH, name="hos.sd.curve")
# legend (kept below the 9 h cap line so it does not mask it)
box(gs, PL + 0.05, HOS_Y + 0.16, 3.95, 0.22, fill="FFFFFF",
    name="hos.legend.bg")
conn(gs, PL + 0.12, HOS_Y + 0.27, PL + 0.38, HOS_Y + 0.27, c_cd, 1.75,
     name="hos.legend.cd.line")
txt(gs, PL + 0.44, HOS_Y + 0.18, 1.75, 0.18,
    [[("consecutive driving  t", 0), ("drv", 30)]], size=9,
    align=PP_ALIGN.LEFT, name="hos.legend.cd.txt")
conn(gs, PL + 2.30, HOS_Y + 0.27, PL + 2.56, HOS_Y + 0.27, c_sd, 1.5,
     dash=MSO_LINE.DASH, name="hos.legend.sd.line")
txt(gs, PL + 2.62, HOS_Y + 0.18, 1.35, 0.18,
    [[("shift driving  t", 0), ("sd", 30)]], size=9,
    align=PP_ALIGN.LEFT, name="hos.legend.sd.txt")
# clock ticks under the HoS panel
th = int(math.ceil(TMIN))
while th <= TMAX:
    if th >= 8:
        conn(gs, xt(th), HOS_Y + HOS_H, xt(th), HOS_Y + HOS_H + 0.04,
             "AAAAAA", 0.75, name=f"hos.xtick.{th}")
        txt(gs, xt(th) - 0.3, HOS_Y + HOS_H + 0.05, 0.6, 0.16,
            f"{th % 24:02d}:00", size=8, color="555555",
            name=f"hos.xticklabel.{th}")
    th += 2

# =========================================================== driver =========
g = slide.shapes.add_group_shape(); g.name = "panel.driver"
gs = g.shapes
txt(gs, 0.10, DRV_Y + DRV_H / 2 - 0.12, 0.85, 0.25, "Driver", size=13,
    bold=True, align=PP_ALIGN.LEFT, name="driver.label")
box(gs, PL, DRV_Y, PW, DRV_H, fill=None, line="AAAAAA", lw=1.0,
    name="driver.frame")
DBY, DBH = DRV_Y + 0.034, DRV_H - 0.068
for (t0, t1) in drive_blocks:
    s = box(gs, xt(t0), DBY, xt(t1) - xt(t0), DBH, fill=c_drive,
            line=c_drive_e, lw=0.5, name=f"driver.drive@{t0:.1f}")
    if t1 - t0 > 1.55:
        s.text_frame.text = "Driving"
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _fmt_run(s.text_frame.paragraphs[0].runs[0], 9, "333333", False, False)

def drv_kind(kind, stop_e):
    if kind in ("mstop", "queue", "mseq", "mlay"):
        return "work"
    if kind == "charge":   # concurrent break credit g, else charging is work
        return "break" if (stop_e["brk"] and not stop_e["sigma"]) else "work"
    if kind == "service":
        return "service"
    return kind

def drv_label(dk, e):
    if dk == "work":
        return "Working time"
    if dk == "service":
        return "Service"
    if dk == "break":
        return {"b45": "45′ break", "b15": "1st split break (15′)",
                "b30": "2nd split break (30′)"}[e["brk"]]
    if dk == "rest":
        return "9 h rest (reduced)" if e["rst"] == "r2" else "11 h rest"
    return None

rot_labels = []
for e in events:
    merged = []
    for kind, t0, t1 in e["blocks"]:
        dk = drv_kind(kind, e)
        if merged and merged[-1][0] == dk:
            merged[-1][2] = t1
        else:
            merged.append([dk, t0, t1])
    for dk, t0, t1 in merged:
        if t1 - t0 < 1e-4:
            continue
        w = xt(t1) - xt(t0)
        if dk == "work":
            box(gs, xt(t0), DBY, w, DBH, fill=c_work, line=c_work_e, lw=0.5,
                name=f"driver.work@{t0:.1f}")
        elif dk == "service":
            box(gs, xt(t0), DBY, w, DBH, fill=c_serv, line=c_serv_e, lw=0.5,
                name=f"driver.service@{t0:.1f}")
        elif dk == "break":
            box(gs, xt(t0), DBY, w, DBH, fill=c_brk_e, back=c_brk_f,
                pattern=MSO_PATTERN.LIGHT_UPWARD_DIAGONAL, line=c_brk_e,
                lw=0.75, name=f"driver.break@{t0:.1f}")
        elif dk == "rest":
            box(gs, xt(t0), DBY, w, DBH, fill=c_rst_e, back=c_rst_f,
                pattern=MSO_PATTERN.WIDE_DOWNWARD_DIAGONAL, line=c_rst_e,
                lw=0.75, name=f"driver.rest@{t0:.1f}")
    for dk, t0, t1 in merged:
        label = drv_label(dk, e)
        if label and t1 - t0 > 1e-4:
            rot_labels.append((0.5 * (t0 + t1), label))

# rotated labels below the driver row (top-right corner anchored under block)
ROT = 38.0
cR, sR = math.cos(math.radians(ROT)), math.sin(math.radians(ROT))
LW_, LH_ = 2.0, 0.24
for xm, label in rot_labels:
    px, py = xt(xm), DRV_Y + DRV_H + 0.06
    ccx = px - (LW_ / 2 * cR - LH_ / 2 * sR)
    ccy = py + (LW_ / 2 * sR + LH_ / 2 * cR)
    tb = txt(gs, ccx - LW_ / 2, ccy - LH_ / 2, LW_, LH_, label, size=9,
             color="555555", align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
             rotation=-ROT, name=f"driver.rotlabel.{label[:12]}@{xm:.1f}")

# =========================================================== caption ========
_cap_src = ("oracle solution of run " if SOURCE == "oracle" else "run ")
txt(slide.shapes, 4.4, 7.22, 8.55, 0.2,
    [[(_cap_src + SOLFILE + "   •   break rules: 45′ after ≤ 4.5 h driving, "
       "splittable 15′+30′; daily rest 11 h (9 h reduced)   •   charge counts "
       "as break when concurrent (g", 0), ("i", -25),
      ("), as work when followed by rest (σ", 0), ("i", -25), (" = 1)", 0)]],
    size=7.5, color="888888", align=PP_ALIGN.RIGHT, name="caption")

prs.save(OUT)
print("saved", OUT)
print("shapes on slide:", len(slide.shapes))
