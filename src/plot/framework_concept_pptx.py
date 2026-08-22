# -*- coding: utf-8 -*-
"""
Editable PowerPoint version of the conceptual FRAMEWORK figure.

Replaces the old "action selector" slide.  The message is the decision
architecture axis: one simulator, one swappable policy slot; the methods
differ only in *when* the decision is committed and *what information* it
uses.  Notation is deliberately kept out -- symbols belong in the model
section, not here.

Every element is a native PowerPoint shape, named by prefix so the Selection
Pane (Alt+F10) is navigable:
    zoneA.* / zoneB.* / zoneC.*   the three commit-time zones
    loop.*                        state -> policy -> drive -> reveal cycle
    shelf.*                       the four policies in the policy slot
    tbl.*                         the bottom comparison strip

Run:  python -m src.plot.framework_concept_pptx
"""
from src import paths as _paths

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

OUT = _paths.figure_out("framework_concept.pptx")

# ============================================================ styling =======
def RGB(h): return RGBColor.from_string(h.lstrip("#"))

INK, MUT, FAINT = "222222", "5F5F5F", "AFAFAF"
FONT  = "Calibri"
EMOJI = "Segoe UI Emoji"

# zone / role palette (kept close to the original slide: orange = uncertainty,
# blue = state, green = decision)
C_OPEN,  F_OPEN  = "B26B00", "FCF2E2"      # open-loop, pre-departure
C_CLOSE, F_CLOSE = "1E7A3C", "E9F4EC"      # closed-loop, online
C_ORAC,  F_ORAC  = "6A4C93", "F1ECF8"      # ex-post bound
C_STATE, F_STATE = "2B6CB0", "E9F1F9"      # state
C_UNC,   F_UNC   = "D2691E", "FDF1E5"      # exogenous information
F_SIM            = "FAFAFA"                # simulator frame
BADGE            = "3F4A56"                # Powell policy-class chip

# ============================================================ layout ========
SLIDE_H = 7.5
prs = Presentation()
prs.slide_width  = Emu(12192000)           # 13.333 in, exact 16:9
prs.slide_height = Inches(SLIDE_H)
sl = prs.slides.add_slide(prs.slide_layouts[6])
sh = sl.shapes

# ============================================================ helpers =======
def box(x, y, w, h, fill=None, line=None, lw=0.75, name=None,
        shape=MSO_SHAPE.RECTANGLE, dash=None):
    s = sh.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = RGB(fill)
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = RGB(line)
        s.line.width = Pt(lw)
        if dash is not None:
            s.line.dash_style = dash
    if name:
        s.name = name
    return s

def _arrowify(c, head=True, tail=False, size="med"):
    """DrawingML arrowheads; headEnd/tailEnd must come last inside <a:ln>."""
    ln = c.line._get_or_add_ln()
    if tail:
        ln.append(ln.makeelement(qn("a:headEnd"),
                                 {"type": "triangle", "w": size, "len": size}))
    if head:
        ln.append(ln.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": size, "len": size}))

def conn(x1, y1, x2, y2, color, lw=1.25, dash=None, name=None,
         head=False, tail=False, size="med"):
    c = sh.add_connector(MSO_CONNECTOR.STRAIGHT,
                         Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.shadow.inherit = False
    c.line.color.rgb = RGB(color)
    c.line.width = Pt(lw)
    if dash is not None:
        c.line.dash_style = dash
    if head or tail:
        _arrowify(c, head=head, tail=tail, size=size)
    if name:
        c.name = name
    return c

def txt(x, y, w, h, lines, size=10, color=INK, bold=False, italic=False,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, wrap=False,
        font=FONT, space=0.0, name=None):
    """lines: str | list of str lines."""
    tb = sh.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for li, ln in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align
        if space and li:
            p.space_before = Pt(space)
        r = p.add_run()
        r.text = ln
        r.font.size   = Pt(size)
        r.font.name   = font
        r.font.bold   = bold
        r.font.italic = italic
        r.font.color.rgb = RGB(color)
    if name:
        tb.name = name
    return tb

def card(x, y, w, h, title, body, accent, fill, name,
         tsize=10.5, bsize=8.0, dash=None):
    box(x, y, w, h, fill=fill, line=accent, lw=1.25, name=name + ".box",
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=dash)
    txt(x + 0.06, y + 0.09, w - 0.12, 0.24, title, size=tsize, bold=True,
        color=accent, name=name + ".title")
    if body:
        txt(x + 0.06, y + 0.09 + 0.24, w - 0.12, h - 0.40,
            body, size=bsize, color=MUT, space=1.5, name=name + ".body")

def pill(x, y, w, h, label, color, name):
    box(x, y, w, h, fill=color, line=color, lw=1.0, name=name + ".box",
        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(x, y + 0.055, w, h, label, size=10.5, bold=True, color="FFFFFF",
        name=name + ".label")

# ======================================================== title ============
txt(0.35, 0.16, 12.63, 0.32,
    "Decision architectures: one simulator, one swappable policy",
    size=17, bold=True, align=PP_ALIGN.LEFT, name="title")
txt(0.35, 0.50, 12.63, 0.24,
    "methods ordered by when the decision is committed, and by what "
    "information it is allowed to use",
    size=10.5, italic=True, color=MUT, align=PP_ALIGN.LEFT, name="subtitle")

# ======================================================== zone frames ======
FY, FH = 1.00, 4.18                                   # frame top / height
AX, AW = 0.35, 3.05                                   # zone A
BX, BW = 3.72, 6.03                                   # zone B
CX, CW = 10.55, 2.43                                  # zone C

box(AX, FY, AW, FH, fill=F_OPEN, line=C_OPEN, lw=1.0, name="zoneA.frame",
    shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=MSO_LINE.DASH)
box(BX, FY, BW, FH, fill=F_SIM, line=C_CLOSE, lw=1.5, name="zoneB.frame",
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
box(CX, FY, CW, FH, fill=F_ORAC, line=C_ORAC, lw=1.0, name="zoneC.frame",
    shape=MSO_SHAPE.ROUNDED_RECTANGLE, dash=MSO_LINE.DASH)

pill(AX + 0.30, FY - 0.17, 2.45, 0.34, "BEFORE DEPARTURE", C_OPEN, "zoneA.pill")
pill(BX + 1.75, FY - 0.17, 2.55, 0.34, "AT EVERY STOP", C_CLOSE, "zoneB.pill")
pill(CX + 0.10, FY - 0.17, 2.23, 0.34, "AFTER THE FACT", C_ORAC, "zoneC.pill")

txt(AX + 0.08, FY + 0.26, AW - 0.16, 0.22, "open-loop  -  no feedback",
    size=9, italic=True, color=C_OPEN, name="zoneA.sub")
txt(BX + 0.08, FY + 0.26, BW - 0.16, 0.22,
    "closed-loop  -  re-decides from the realized state",
    size=9, italic=True, color=C_CLOSE, name="zoneB.sub")
txt(CX + 0.08, FY + 0.26, CW - 0.16, 0.22, "not a policy  -  a bound",
    size=9, italic=True, color=C_ORAC, name="zoneC.sub")

# ======================================================== zone A: plans ====
card(0.50, 1.66, 2.40, 1.28, "RO  -  robust plan",
     ["worst case over an interval",
      "uncertainty set",
      "whole schedule fixed at departure,",
      "then executed as-is"],
     C_OPEN, "FFFFFF", "zoneA.ro")

card(0.50, 3.16, 2.40, 1.52, "2SP  -  two-stage plan",
     ["expected arrival over a scenario set",
      "activity structure fixed at departure;",
      "durations re-optimised online, with a",
      "repair step if the plan breaks"],
     C_OPEN, "FFFFFF", "zoneA.2sp")

# overhead lane: precomputed plans enter the policy slot without seeing state
LANE_Y, RISER_X, LANE_END = 1.78, 3.14, 8.10
conn(2.90, 2.30, RISER_X, 2.30, C_OPEN, 1.25, name="zoneA.stub.ro")
conn(2.90, 3.92, RISER_X, 3.92, C_OPEN, 1.25, name="zoneA.stub.2sp")
conn(RISER_X, 3.92, RISER_X, LANE_Y, C_OPEN, 1.25, name="zoneA.riser")
conn(RISER_X, LANE_Y, LANE_END, LANE_Y, C_OPEN, 1.25, name="zoneA.lane")
conn(LANE_END, LANE_Y, LANE_END, 1.99, C_OPEN, 1.25, head=True,
     name="zoneA.lane.drop")
txt(4.05, LANE_Y - 0.28, 3.90, 0.22,
    "computed once, before any travel time is known",
    size=8.5, italic=True, color=C_OPEN, name="zoneA.lane.label")

# ======================================================== zone B: the loop =
STX, STY, STW, STH = 3.98, 2.18, 1.95, 0.98           # state
EXX, EXY, EXW, EXH = 3.98, 3.70, 1.95, 1.15           # drive + reveal
PLX, PLY, PLW, PLH = 6.35, 1.95, 3.25, 2.97           # policy slot

box(STX, STY, STW, STH, fill=F_STATE, line=C_STATE, lw=1.25,
    name="loop.state.box", shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(STX, STY + 0.10, STW, 0.24, "State at stop i", size=11, bold=True,
    color=C_STATE, name="loop.state.title")
txt(STX + 0.08, STY + 0.42, STW - 0.16, 0.55,
    ["time  -  charge left",
     "driving since last break",
     "time since last rest"],
    size=8.5, color=MUT, name="loop.state.body")

box(EXX, EXY, EXW, EXH, fill=F_UNC, line=C_UNC, lw=1.25,
    name="loop.exec.box", shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(EXX, EXY + 0.09, EXW, 0.24, "Drive the leg", size=11, bold=True,
    color=C_UNC, name="loop.exec.title")
txt(EXX, EXY + 0.38, EXW, 0.28, "⏱  ⚡", size=13, color=C_UNC,
    font=EMOJI, name="loop.exec.icons")
txt(EXX + 0.06, EXY + 0.72, EXW - 0.12, 0.44,
    ["travel time and energy use", "are revealed"],
    size=8.5, color=MUT, name="loop.exec.body")

# policy slot
box(PLX, PLY, PLW, PLH, fill="FFFFFF", line=C_CLOSE, lw=2.0,
    name="shelf.box", shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(PLX, PLY + 0.10, PLW, 0.26, "POLICY", size=12.5, bold=True, color=C_CLOSE,
    name="shelf.title")
txt(PLX, PLY + 0.37, PLW, 0.22, "chooses the action at this stop",
    size=8.5, italic=True, color=MUT, name="shelf.sub")

SHELF = [
    ("RO",     "CFA", "executes the committed plan",       C_OPEN,  F_OPEN),
    ("2SP",    "DLA", "structure fixed, LP on durations",  C_OPEN,  F_OPEN),
    ("GREEDY", "PFA", "priority rule, no optimisation",    C_CLOSE, F_CLOSE),
    ("LA",     "DLA", "lookahead, re-solved at each stop", C_CLOSE, F_CLOSE),
]
CY, CH, CG, BW_ = PLY + 0.61, 0.48, 0.05, 0.55
for k, (nm, cls, desc, acc, fil) in enumerate(SHELF):
    y  = CY + k * (CH + CG)
    bx = PLX + PLW - 0.13 - BW_
    low = nm.lower()
    box(PLX + 0.13, y, PLW - 0.26, CH, fill=fil, line=acc, lw=1.0,
        name="shelf.%s.box" % low, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(PLX + 0.22, y + 0.05, 0.72, 0.22, nm, size=10, bold=True, color=acc,
        align=PP_ALIGN.LEFT, name="shelf.%s.name" % low)
    txt(PLX + 0.94, y + 0.07, bx - PLX - 1.00, 0.22, desc, size=8, color=MUT,
        align=PP_ALIGN.LEFT, name="shelf.%s.desc" % low)
    box(bx, y + 0.11, BW_, 0.26, fill=BADGE, line=BADGE, lw=1.0,
        name="shelf.%s.class.box" % low, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(bx, y + 0.155, BW_, 0.20, cls, size=8, bold=True, color="FFFFFF",
        name="shelf.%s.class" % low)

# cycle arrows
conn(STX + STW, 2.45, PLX, 2.45, C_STATE, 1.5, head=True, name="loop.arr.state")
txt(STX + STW - 0.02, 2.18, PLX - STX - STW + 0.04, 0.20, "state",
    size=8.5, italic=True, color=C_STATE, name="loop.arr.state.label")

conn(PLX, 4.76, STX + STW, 4.76, C_CLOSE, 1.5, head=True, name="loop.arr.act")
txt(STX + STW - 0.02, 4.49, PLX - STX - STW + 0.04, 0.20, "action",
    size=8.5, italic=True, color=C_CLOSE, name="loop.arr.act.label")

conn(STX + STW / 2, EXY, STX + STW / 2, STY + STH, MUT, 1.5, head=True,
     name="loop.arr.next")
txt(STX + STW / 2 + 0.07, STY + STH + 0.12, 1.30, 0.20, "next stop",
    size=8.5, italic=True, color=MUT, align=PP_ALIGN.LEFT,
    name="loop.arr.next.label")

txt(BX + 0.08, FY + FH - 0.23, BW - 0.16, 0.22,
    "one discrete-event simulator  -  identical sample paths for every method",
    size=8.5, italic=True, color=FAINT, name="zoneB.footer")

# ======================================================== zone C: oracle ===
card(CX + 0.16, 1.66, CW - 0.32, 1.62, "ORACLE",
     ["hindsight optimum: the full-route",
      "problem re-solved once every",
      "travel time is known",
      "",
      "the best any method could have done"],
     C_ORAC, "FFFFFF", "zoneC.oracle", tsize=12, dash=MSO_LINE.DASH)

txt(CX + 0.20, 3.44, CW - 0.40, 0.80,
    ["not a policy: it uses information",
     "no driver could have had",
     "(it violates nonanticipativity)"],
    size=8.5, italic=True, color=C_ORAC, name="zoneC.note")

conn(BX + BW, 4.58, CX, 4.58, C_ORAC, 1.5, head=True, tail=True,
     name="zoneC.gap.arrow")
txt(BX + BW - 0.06, 4.28, CX - BX - BW + 0.12, 0.22, "gap Δ",
    size=9.5, bold=True, color=C_ORAC, name="zoneC.gap.label")

# ======================================================== comparison strip =
TY = 5.42
LBX, LBW = 0.35, 1.42
COLW, COLG = 2.12, 0.10
COLX = [1.85 + k * (COLW + COLG) for k in range(5)]
COLS = [
    ("RO",     C_OPEN,  F_OPEN,
     ["the whole schedule,", "at departure"],
     ["an interval uncertainty set", "(no realizations)"],
     ["CFA - worst-case parameters"]),
    ("2SP",    C_OPEN,  F_OPEN,
     ["structure at departure,", "durations at each stop"],
     ["a scenario set,", "then the realized state"],
     ["DLA - two-stage, solved once"]),
    ("GREEDY", C_CLOSE, F_CLOSE,
     ["one action,", "at each stop"],
     ["the current state only"],
     ["PFA - analytic decision rule"]),
    ("LA",     C_CLOSE, F_CLOSE,
     ["one action,", "at each stop"],
     ["current state + scenarios", "over a horizon L"],
     ["DLA - rolling-horizon lookahead"]),
    ("ORACLE", C_ORAC,  F_ORAC,
     ["nothing - computed", "after the run"],
     ["the entire realized path"],
     ["wait-and-see bound (no class)"]),
]
ROWS = [("commits", 5.80, 0.44), ("information used", 6.28, 0.44),
        ("policy class", 6.76, 0.36)]

for lab, y, h in ROWS:
    txt(LBX, y + 0.10, LBW, 0.22, lab, size=9, bold=True, color=MUT,
        align=PP_ALIGN.RIGHT, name="tbl.rowlabel." + lab.split()[0])

for k, (nm, acc, fil, commits, info, cls) in enumerate(COLS):
    x = COLX[k]
    box(x, TY, COLW, 0.30, fill=acc, line=acc, lw=1.0,
        name="tbl.%s.head" % nm.lower(), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(x, TY + 0.055, COLW, 0.24, nm, size=10, bold=True, color="FFFFFF",
        name="tbl.%s.name" % nm.lower())
    for (lab, y, h), val in zip(ROWS, (commits, info, cls)):
        last = lab == "policy class"
        box(x, y, COLW, h, fill=fil if last else "FFFFFF", line=FAINT,
            lw=0.5, name="tbl.%s.%s" % (nm.lower(), lab.split()[0]))
        txt(x + 0.05, y + (0.10 if last else 0.08), COLW - 0.10, h - 0.10,
            val, size=8, bold=last, color=acc if last else INK,
            name="tbl.%s.%s.txt" % (nm.lower(), lab.split()[0]))

LEG = [
    ("PFA", "an analytic decision rule",          INK),
    ("CFA", "a hedged deterministic model",       INK),
    ("DLA", "an approximate model of the future", INK),
    ("VFA", "not used here",                      FAINT),
]
LGY = 7.16
txt(LBX, LGY - 0.02, LBW, 0.22, "policy classes", size=8.5, bold=True,
    color=MUT, align=PP_ALIGN.RIGHT, name="tbl.legend.rowlabel")
txt(LBX, LGY + 0.16, LBW, 0.20, "Powell (2022)", size=7.5, italic=True,
    color=FAINT, align=PP_ALIGN.RIGHT, name="tbl.legend.cite")
for k, (ac, desc, col) in enumerate(LEG):
    x = 1.85 + k * 2.75
    c = BADGE if col == INK else FAINT
    box(x, LGY, 0.45, 0.24, fill=c, line=c, lw=1.0,
        name="tbl.legend.%s.box" % ac.lower(),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(x, LGY + 0.035, 0.45, 0.20, ac, size=8, bold=True, color="FFFFFF",
        name="tbl.legend.%s.ac" % ac.lower())
    txt(x + 0.53, LGY + 0.05, 2.15, 0.20, desc, size=8, color=col,
        align=PP_ALIGN.LEFT, name="tbl.legend.%s.txt" % ac.lower())

prs.save(OUT)
print("saved", OUT)
